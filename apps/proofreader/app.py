import io
import os
from pathlib import Path

from flask import Flask, jsonify, render_template, request
import anthropic

app = Flask(__name__)

ALLOWED = {'.pdf', '.pptx', '.docx'}

SYSTEM_PROMPT = """\
# 役割
あなたは研修資料の校正専門エージェントです。
ユースフル株式会社が企業クライアントの事務局・受講者に提出する研修資料（PDF）を提出前にチェックし、
信頼性の高い資料として仕上げることをサポートします。

# タスク
アップロードされたPDFを精読し、以下の4カテゴリのミスを検出して、表形式でレポートしてください。

## チェックカテゴリ
1. **表記ゆれ**：同じ概念・語句が複数の表記で使われている（例：「パソコン」「PC」「ＰＣ」）
2. **誤字・脱字**：明らかな誤字、脱字、余分な文字
3. **句読点・記号の誤り**：不自然な句読点、全角/半角の混在、括弧の対応ミスなど
4. **敬語・文体の不統一**：敬語レベルの混在（です/ます体とだ/である体の混在）、受講者/事務局向けとして不適切な表現

# 出力形式
以下の表形式で出力してください。

| 優先度 | ページ | カテゴリ | 該当箇所 | 問題の内容 | 修正案 |
|--------|--------|----------|----------|------------|--------|

## 優先度の基準（⭐で表示）
- ⭐⭐⭐：意味が変わる・読者に誤解を与える・信頼性を大きく損なうミス
- ⭐⭐　：読みにくさや不統一が目立つが意味は通じるミス
- ⭐　　：細かい表記の揺れや軽微な統一性の問題

# 出力の流れ
1. **サマリー**：検出したミスの総数をカテゴリ別に1〜2行で述べる
2. **校正レポート表**：上記の表形式で全件出力
3. **総評**：資料全体の品質に関する一言コメント（提出可否の目安も含める）

# 注意事項
- 固有名詞・商品名・略語は文脈から判断し、誤検知に注意すること
- Microsoft 365関連の用語（例：Copilot、Teams、SharePointなど）は正式表記を基準にすること
- 「ページ」が特定できない場合はスライド番号や見出し名で代替すること
- 検出されなかったカテゴリは「該当なし」と明記すること
"""


def extract_text_from_pdf(file_bytes: bytes) -> str:
    import pdfplumber
    parts = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text() or ""
            if text.strip():
                parts.append(f"=== ページ {i} ===\n{text}")
    return "\n\n".join(parts)


def extract_text_from_pptx(file_bytes: bytes) -> str:
    from pptx import Presentation
    parts = []
    prs = Presentation(io.BytesIO(file_bytes))
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            parts.append(f"=== スライド {i} ===\n" + "\n".join(texts))
    return "\n\n".join(parts)


def extract_text_from_docx(file_bytes: bytes) -> str:
    import docx
    doc = docx.Document(io.BytesIO(file_bytes))
    return "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())


def extract_text(file_bytes: bytes, ext: str) -> str:
    if ext == '.pdf':
        return extract_text_from_pdf(file_bytes)
    elif ext == '.pptx':
        return extract_text_from_pptx(file_bytes)
    elif ext == '.docx':
        return extract_text_from_docx(file_bytes)
    raise ValueError(f"非対応: {ext}")


@app.route('/')
def index():
    return render_template('index.html')


@app.route('/api-status')
def api_status():
    return jsonify({'has_key': bool(os.environ.get('ANTHROPIC_API_KEY'))})


@app.route('/proofread', methods=['POST'])
def proofread():
    if 'file' not in request.files:
        return jsonify({'error': 'ファイルがありません'}), 400

    file = request.files['file']
    if not file.filename:
        return jsonify({'error': 'ファイル名が空です'}), 400

    ext = Path(file.filename).suffix.lower()
    if ext not in ALLOWED:
        return jsonify({'error': f'非対応の形式です: {ext}（PDF / PPTX / DOCX のみ対応）'}), 400

    api_key = request.form.get('api_key') or os.environ.get('ANTHROPIC_API_KEY')
    if not api_key:
        return jsonify({'error': 'Anthropic API キーが設定されていません'}), 400

    file_bytes = file.read()
    if len(file_bytes) > 20 * 1024 * 1024:
        return jsonify({'error': 'ファイルサイズが 20MB を超えています'}), 400

    try:
        text = extract_text(file_bytes, ext)
    except Exception as e:
        return jsonify({'error': f'テキスト抽出エラー: {e}'}), 500

    if not text.strip():
        return jsonify({'error': 'テキストを抽出できませんでした（画像PDFは非対応です）'}), 400

    client = anthropic.Anthropic(api_key=api_key)
    try:
        msg = client.messages.create(
            model='claude-opus-4-6',
            max_tokens=4096,
            system=SYSTEM_PROMPT,
            messages=[{'role': 'user', 'content': f'以下の資料を校正してください。\n\n{text}'}],
        )
        return jsonify({'result': msg.content[0].text})
    except anthropic.AuthenticationError:
        return jsonify({'error': 'API キーが無効です。正しいキーを入力してください。'}), 401
    except Exception as e:
        return jsonify({'error': f'API エラー: {e}'}), 500


if __name__ == '__main__':
    app.run(port=5002, debug=False)
